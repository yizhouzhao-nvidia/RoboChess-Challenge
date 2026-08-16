from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "RoboChess_IsaacLab_CoRL2026_Workshop.pptx"
RENDERS = ROOT / "lab" / "images" / "renders"

NAVY = RGBColor(13, 22, 42)
BLUE = RGBColor(58, 116, 255)
CYAN = RGBColor(64, 210, 214)
INK = RGBColor(21, 31, 48)
MUTED = RGBColor(92, 104, 121)
PALE = RGBColor(242, 246, 252)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(42, 174, 120)
ORANGE = RGBColor(245, 151, 65)


def add_text(slide, text, x, y, w, h, size=20, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def rect(slide, x, y, w, h, fill, radius=False, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill if line is None else line
    return shape


def title(slide, heading, sub=""):
    add_text(slide, heading, 0.7, 0.42, 11.8, 0.55, 28, NAVY, True)
    if sub:
        add_text(slide, sub, 0.73, 0.98, 11.4, 0.35, 11, MUTED)
    rect(slide, 0.7, 1.28, 1.2, 0.055, BLUE)


def footer(slide, number):
    add_text(slide, "RoboChess × Isaac Lab", 0.72, 7.14, 3, 0.18, 9, MUTED)
    add_text(slide, f"{number} / 5", 11.65, 7.14, 0.6, 0.18, 9, MUTED, align=PP_ALIGN.RIGHT)


def bullet(slide, text, x, y, w, accent=BLUE, size=17):
    rect(slide, x, y + 0.12, 0.1, 0.1, accent, radius=True)
    add_text(slide, text, x + 0.22, y, w - 0.22, 0.5, size, INK)


def card(slide, heading, body, x, y, w, h, accent=BLUE):
    rect(slide, x, y, w, h, WHITE, radius=True, line=RGBColor(221, 229, 240))
    rect(slide, x, y, 0.08, h, accent)
    add_text(slide, heading, x + 0.28, y + 0.22, w - 0.5, 0.34, 16, INK, True)
    add_text(slide, body, x + 0.28, y + 0.66, w - 0.5, h - 0.82, 12, MUTED)


def img(slide, filename, x, y, w, h):
    path = RENDERS / filename
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    rect(slide, x, y, w, h, WHITE, line=RGBColor(205, 215, 230))


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 — title
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, 13.333, 7.5, NAVY)
    rect(s, 0, 0, 13.333, 0.12, CYAN)
    add_text(s, "RoboChess", 0.8, 1.08, 7.4, 0.7, 44, WHITE, True)
    add_text(s, "A configurable Isaac Lab task for sim-to-real robot learning", 0.84, 1.82, 7.6, 0.5, 22, RGBColor(202, 216, 242))
    add_text(s, "CoRL 2026 Workshop • Sim-Real Robot Learning", 0.84, 2.55, 6.8, 0.35, 14, CYAN, True)
    add_text(s, "From board configuration to embodiment diversity — a single-arm visual task ready for policy development.", 0.84, 3.08, 6.1, 0.75, 16, WHITE)
    img(s, "ur10_1d.png", 8.05, 0.72, 4.45, 3.22)
    img(s, "so101_1d.png", 8.05, 4.14, 2.15, 1.56)
    img(s, "piper_1d.png", 10.35, 4.14, 2.15, 1.56)
    add_text(s, "Implementation snapshot • August 2026", 0.84, 6.64, 4.5, 0.25, 11, RGBColor(184, 199, 227))

    # 2 — workshop motivation
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, 13.333, 7.5, PALE)
    title(s, "Why RoboChess? A compact testbed for the sim-to-real gap", "A challenge format makes perception, manipulation, planning, and robustness measurable together.")
    card(s, "Perceive", "Board state, piece identity, pose, and camera variation.", 0.8, 1.72, 3.85, 1.48, CYAN)
    card(s, "Manipulate", "Pick, place, avoid collisions, and execute legal moves.", 4.75, 1.72, 3.85, 1.48, BLUE)
    card(s, "Generalize", "Across simulators, robot embodiments, and real-world perturbations.", 8.7, 1.72, 3.85, 1.48, ORANGE)
    add_text(s, "Workshop research questions", 0.82, 3.72, 4.4, 0.35, 18, NAVY, True)
    bullet(s, "Uniform, quantitative evaluation of robot-learning methods", 0.88, 4.22, 5.6, BLUE, 15)
    bullet(s, "From visually plausible rollouts to repeatable long-horizon task solving", 0.88, 4.84, 5.8, BLUE, 15)
    bullet(s, "Robustness under lighting, viewpoint, object-pose, and embodiment shifts", 0.88, 5.46, 5.8, BLUE, 15)
    rect(s, 7.25, 3.74, 5.22, 2.2, NAVY, radius=True)
    add_text(s, "RoboChess contribution", 7.65, 4.08, 4.4, 0.35, 19, WHITE, True)
    add_text(s, "A standardized sim-to-real manipulation challenge that combines chess reasoning with physical execution.", 7.65, 4.62, 4.15, 0.76, 15, RGBColor(218, 229, 248))
    add_text(s, "Source: CoRL 2026 workshop proposal and challenge site", 0.82, 6.65, 6.5, 0.2, 10, MUTED)
    footer(s, 2)

    # 3 — task composition
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, 13.333, 7.5, WHITE)
    title(s, "Isaac Lab task: one interface, four chess scales, six embodiments", "Manager-based scene configuration keeps the first visual stage simple and composable.")
    add_text(s, "Chess scenario argument", 0.8, 1.65, 3.1, 0.35, 18, NAVY, True)
    scenarios = [("1D", "six-piece line", CYAN), ("3 × 3", "pawn microgame", GREEN), ("4 × 4", "minichess", BLUE), ("8 × 8", "standard chess", ORANGE)]
    for i, (name, sub, c) in enumerate(scenarios):
        x = 0.82 + i * 1.53
        rect(s, x, 2.15, 1.34, 1.05, c, radius=True)
        add_text(s, name, x, 2.36, 1.34, 0.28, 19, WHITE, True, PP_ALIGN.CENTER)
        add_text(s, sub, x, 2.72, 1.34, 0.2, 9, WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "Robot argument", 7.25, 1.65, 2.6, 0.35, 18, NAVY, True)
    robots = ["SO-101", "Piper", "UR10", "Flexiv Rizon", "Rebot", "YAM"]
    for i, name in enumerate(robots):
        x = 7.26 + (i % 3) * 1.67
        y = 2.15 + (i // 3) * 0.72
        rect(s, x, y, 1.47, 0.5, PALE, radius=True, line=RGBColor(218, 226, 238))
        add_text(s, name, x, y + 0.13, 1.47, 0.2, 12, INK, True, PP_ALIGN.CENTER)
    rect(s, 0.82, 4.0, 11.66, 1.48, PALE, radius=True)
    add_text(s, "Launch contract", 1.12, 4.25, 2.0, 0.28, 16, NAVY, True)
    code = "python scripts/zero_agent.py --task RoboChess-UR10-Visual-v0 --chess_scenario 1d --robot ur10"
    rect(s, 3.18, 4.22, 8.72, 0.5, NAVY, radius=True)
    add_text(s, code, 3.45, 4.36, 8.18, 0.18, 12, WHITE)
    add_text(s, "Local chess USD assets • table and camera • stationary single-arm visual scene • no teleoperation", 1.12, 4.98, 10.5, 0.22, 12, MUTED)
    footer(s, 3)

    # 4 — visual validation
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, 13.333, 7.5, PALE)
    title(s, "Visual validation: 1D chess rendered across robot options", "The corrected camera explicitly targets the tabletop, chess pieces, and robot workspace.")
    names = [("so101_1d.png", "SO-101"), ("piper_1d.png", "Piper"), ("ur10_1d.png", "UR10"),
             ("flexiv_rizon_1d.png", "Flexiv Rizon"), ("rebot_1d.png", "Rebot"), ("yam_1d.png", "YAM")]
    for i, (file, label) in enumerate(names):
        col, row = i % 3, i // 3
        x, y = 0.75 + col * 4.15, 1.65 + row * 2.42
        img(s, file, x, y, 3.73, 1.82)
        add_text(s, label, x + 0.03, y + 1.94, 3.65, 0.2, 12, INK, True, PP_ALIGN.CENTER)
    add_text(s, "All six configurations were instantiated for the same 1D board layout. These are scene-validation renders, not policy results.", 0.82, 6.62, 11.6, 0.24, 11, MUTED)
    footer(s, 4)

    # 5 — next milestones
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, 13.333, 7.5, NAVY)
    title_color = WHITE
    add_text(s, "From visual task to benchmark-ready learning loop", 0.72, 0.44, 11.8, 0.52, 28, title_color, True)
    add_text(s, "A staged path from reliable scene construction to reproducible sim-to-real evaluation.", 0.74, 1.0, 10.8, 0.28, 12, RGBColor(194, 210, 239))
    rect(s, 0.72, 1.35, 1.2, 0.055, CYAN)
    stages = [("01", "Task mechanics", "Piece poses, grasp logic, legal actions, and resets", CYAN),
              ("02", "Learning interface", "Observations, rewards, success metrics, and baselines", BLUE),
              ("03", "Robustness suite", "Camera, lighting, pose, and embodiment randomization", ORANGE),
              ("04", "Sim-to-real track", "Calibration, demonstrations, real robot evaluation", GREEN)]
    for i, (num, head, body, c) in enumerate(stages):
        x = 0.72 + i * 3.12
        rect(s, x, 2.0, 2.72, 2.45, RGBColor(25, 39, 69), radius=True, line=RGBColor(59, 78, 117))
        add_text(s, num, x + 0.28, 2.28, 0.5, 0.24, 13, c, True)
        add_text(s, head, x + 0.28, 2.74, 2.15, 0.36, 16, WHITE, True)
        add_text(s, body, x + 0.28, 3.34, 2.15, 0.66, 12, RGBColor(202, 214, 236))
    rect(s, 0.72, 5.28, 11.92, 0.88, RGBColor(25, 39, 69), radius=True, line=RGBColor(59, 78, 117))
    add_text(s, "Current deliverable", 1.02, 5.55, 2.2, 0.24, 14, CYAN, True)
    add_text(s, "A versioned Isaac Lab RoboChess scene with scenario and robot selection, local chess assets, and camera-validated renders.", 3.05, 5.53, 8.7, 0.28, 13, WHITE)
    add_text(s, "References: CoRL 2026 Workshop Proposal • sites.google.com/view/corl-workshop-robochess/home", 0.74, 6.79, 11.8, 0.2, 10, RGBColor(177, 196, 228))
    add_text(s, "5 / 5", 11.65, 7.14, 0.6, 0.18, 9, RGBColor(177, 196, 228), align=PP_ALIGN.RIGHT)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
